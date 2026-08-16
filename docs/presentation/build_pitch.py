#!/usr/bin/env python3
"""Build eVault hackathon pitch deck."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

OUT = Path(__file__).resolve().parent / "eVault_Pitch.pptx"

# Palette — institutional navy / teal (not purple AI-default)
NAVY = RGBColor(0x0B, 0x1F, 0x33)
TEAL = RGBColor(0x0F, 0x6E, 0x56)
ACCENT = RGBColor(0xC4, 0x5C, 0x26)
LIGHT = RGBColor(0xF4, 0xF7, 0xF5)
MUTED = RGBColor(0x5A, 0x6B, 0x73)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xE8, 0xEF, 0xEC)


def set_run(run, size=18, bold=False, color=NAVY, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color=LIGHT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def title_text(slide, text, top=0.35, size=32):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.2), Inches(0.8))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=True, color=NAVY)


def body_bullets(slide, items, top=1.3, left=0.55, width=12.2, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = item
        set_run(run, size=size, color=NAVY)


def footer(slide, page, total=12):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"eVault  ·  Blockchain Legal Records  ·  {page}/{total}"
    set_run(run, size=11, color=MUTED)


def card(slide, left, top, w, h, title, body):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.fill.background()
    shape.adjustments[0] = 0.08
    tb = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.2), Inches(w - 0.4), Inches(h - 0.35)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run(r, size=16, bold=True, color=TEAL)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    r2 = p2.add_run()
    r2.text = body
    set_run(r2, size=13, color=NAVY)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Title
    s = prs.slides.add_slide(blank)
    add_bg(s, NAVY)
    accent = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.8), Inches(13.333), Inches(1.7)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "eVault"
    set_run(r, size=54, bold=True, color=WHITE)
    box2 = s.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.2))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    r = p.add_run()
    r.text = "Blockchain-based eVault for Legal Records"
    set_run(r, size=26, color=WHITE)
    p2 = tf2.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Secure · Transparent · Accessible for Lawyers, Judges & Clients"
    set_run(r2, size=16, color=RGBColor(0xC8, 0xE6, 0xDC))
    box3 = s.shapes.add_textbox(Inches(0.8), Inches(6.15), Inches(11.5), Inches(0.9))
    tf3 = box3.text_frame
    p = tf3.paragraphs[0]
    r = p.add_run()
    r.text = "Hackathon Prototype  |  Design Document  |  Business Plan"
    set_run(r, size=16, bold=True, color=WHITE)
    p2 = tf3.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Aligned to Ministry of Law & Justice — Access to Justice"
    set_run(r2, size=13, color=WHITE)

    # 2 Problem
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_text(s, "The Problem")
    body_bullets(
        s,
        [
            "Fragmented paper and siloed digital case files slow proceedings",
            "Informal sharing (email, USB) risks confidentiality breaches",
            "Hard to prove which copy is authentic — weak chain of custody",
            "Litigants struggle to access their own authorized records",
            "Result: delays, higher cost, and reduced trust in the system",
        ],
    )
    footer(s, 2)

    # 3 Solution
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_text(s, "Our Solution — eVault")
    body_bullets(
        s,
        [
            "Encrypt legal PDFs with AES-256-GCM before storage",
            "Store ciphertext on IPFS; metadata in secure databases",
            "Register CID + integrity hash on Ethereum (Sepolia prototype)",
            "Role-based apps for lawyers, judges, and citizens",
            "Wallet/JWT auth, explicit share grants, full audit trail",
        ],
        top=1.2,
    )
    card(
        s,
        0.55,
        5.2,
        12.2,
        1.4,
        "Design principle",
        "Confidentiality off-chain  ·  Verifiability on-chain  ·  Accessibility by role. PDFs never go on the blockchain.",
    )
    footer(s, 3)

    # 4 Architecture
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_text(s, "System Architecture")
    cards = [
        (0.5, "Frontend", "React UI\nLawyer · Judge · Citizen"),
        (3.5, "Gateway", "Spring Cloud\nSingle API entry :8080"),
        (6.5, "Services", "Auth · Docs · Chain\nAudit · Notify · Integrate"),
        (9.5, "Anchors", "MySQL · IPFS\nEthereum Sepolia"),
    ]
    for left, t, b in cards:
        card(s, left, 1.4, 2.8, 2.2, t, b)
    body_bullets(
        s,
        [
            "Microservices scale and upgrade independently",
            "Gateway routes /api/documents, /blockchain, /audit, /ecourts, …",
            "Smart contract eVault.sol manages store, verify, roles, permissions",
        ],
        top=4.0,
        size=16,
    )
    footer(s, 4)

    # 5 Security
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_text(s, "Privacy & Security")
    items = [
        ("Encryption", "AES-256-GCM + HKDF; versioned payloads; master key in env only"),
        ("Access control", "JWT + wallet; owner or time-bounded share grant required to download"),
        ("Integrity", "SHA-256 hash; MySQL vs Sepolia verify without decrypting"),
        ("Accountability", "Audit events: upload, share, download, deny, verify"),
    ]
    for i, (t, b) in enumerate(items):
        row, col = divmod(i, 2)
        card(s, 0.5 + col * 6.3, 1.35 + row * 2.4, 6.0, 2.1, t, b)
    footer(s, 5)

    # 6 Smart contracts
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_text(s, "Blockchain Layer — Ethereum Sepolia")
    body_bullets(
        s,
        [
            "Smart contract: eVault.sol (Solidity 0.8.x)",
            "Stores: docId, caseId, IPFS CID, docType, uploader, timestamp, status",
            "Optional content hash for stronger integrity checks",
            "Roles: CLIENT, LAWYER, JUDGE, ADMIN — assignable on-chain",
            "Events for stored, shared, revoked, signed documents",
            "Node service calls the contract; UI never needs users to pay gas for every view",
        ],
    )
    footer(s, 6)

    # 7 UIs
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_text(s, "Stakeholder Interfaces")
    for left, t, b in [
        (0.5, "Lawyer Filing", "Secure & register documents\nGrant access\nDocument receipt"),
        (4.7, "Judge Workspace", "Sign & register judicial orders\nCase-linked filing"),
        (8.9, "Citizen Vault", "Load by case ID\nVerify integrity\nAuthorized download"),
    ]:
        card(s, left, 1.5, 3.9, 3.5, t, b)
    footer(s, 7)

    # 8 Integration
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_text(s, "Interoperability")
    body_bullets(
        s,
        [
            "Integration service: case registry APIs (eCourts-style)",
            "Document classifier assists filing type selection",
            "Identity commitment binds wallet without storing raw Aadhaar in the vault",
            "Gateway APIs ready for CMS / DigiLocker-class connectors in later phases",
            "Prototype proves the path — not a fake hard-coded case list",
        ],
    )
    footer(s, 8)

    # 9 Demo
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_text(s, "Live Demo Flow")
    body_bullets(
        s,
        [
            "1. Connect MetaMask — real JWT auth",
            "2. Lawyer: upload PDF → encrypt → IPFS → chain → MySQL → audit",
            "3. Show document receipt (Doc ID, CID, status)",
            "4. Verify integrity — MySQL vs Sepolia, no decrypt",
            "5. Citizen/owner: authorized DOWNLOAD DOCUMENT",
            "6. Optional: grant access to another wallet; Judge order path",
        ],
    )
    footer(s, 9)

    # 10 Business
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_text(s, "Impact & Business Model")
    card(
        s,
        0.5,
        1.3,
        6.0,
        4.8,
        "Impact",
        "Faster retrieval of authentic filings\n\nLower cost vs paper/email sprawl\n\nHigher trust via integrity proofs\n\nBetter access to justice for clients\n\nTransparent sharing with audit",
    )
    card(
        s,
        6.8,
        1.3,
        6.0,
        4.8,
        "Revenue (ethical)",
        "Court / state platform licenses\n\nIntegration & implementation fees\n\nAnnual maintenance (AMC)\n\nLaw-firm / LSA subscriptions\n\nNever monetize confidential case content",
    )
    footer(s, 10)

    # 11 Roadmap
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_bar(s)
    title_text(s, "Roadmap")
    for i, (t, b) in enumerate(
        [
            ("P0 Prototype", "Working vault + UIs\nSepolia + encryption\n(Hackathon now)"),
            ("P1 Pilot", "One court complex\nTraining + funded ops\n90-day KPIs"),
            ("P2 Integrate", "Official eCourts/CMS\nCitizen retrieval policy"),
            ("P3 Scale", "Multi-state HA\nSecurity audit\nL2 / permissioned options"),
        ]
    ):
        card(s, 0.45 + i * 3.2, 1.6, 3.0, 4.0, t, b)
    footer(s, 11)

    # 12 Close
    s = prs.slides.add_slide(blank)
    add_bg(s, NAVY)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Deliverables"
    set_run(r, size=20, bold=True, color=TEAL)
    for line in [
        "1. Functional prototype — upload, retrieve, share, verify",
        "2. Detailed design document — docs/DESIGN_DOCUMENT.md",
        "3. Business plan — docs/BUSINESS_PLAN.md",
        "4. This presentation + live demo",
    ]:
        p = tf.add_paragraph()
        p.space_before = Pt(8)
        r = p.add_run()
        r.text = line
        set_run(r, size=18, color=WHITE)
    box2 = s.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0))
    tf2 = box2.text_frame
    p = tf2.paragraphs[0]
    r = p.add_run()
    r.text = "Thank you — questions welcome"
    set_run(r, size=28, bold=True, color=WHITE)
    p2 = tf2.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Confidentiality off-chain · Verifiability on-chain · Access by role"
    set_run(r2, size=14, color=RGBColor(0xC8, 0xE6, 0xDC))

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
